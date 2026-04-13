#!/usr/bin/env python3
"""
Pitch Deck Generator by UniqueClub
Generates a professional 10-page VC-ready pitch deck as a .pptx file.

Usage:
    python generate_pitch_deck.py --input startup_data.json --output MyStartup_BP.pptx

Input JSON structure:
{
    "project_name": "My Startup",
    "tagline": "One-line positioning",
    "value_props": ["Point 1", "Point 2", "Point 3"],
    "pain_points": [{"title": "", "description": ""}],
    "market_size": {"tam": "", "sam": "", "som": ""},
    "solution": {"features": [], "differentiation": ""},
    "business_model": {"revenue_streams": [], "pricing": ""},
    "product_status": {"demo": true, "users": "", "metrics": ""},
    "competitors": [{"name": "", "comparison": ""}],
    "team": [{"name": "", "role": "", "background": ""}],
    "traction": {"milestones": [], "metrics": ""},
    "roadmap": {"short": [], "mid": [], "long": []},
    "fundraising": {"amount": "", "equity": "", "use_of_funds": []},
    "language": "zh",
    "colors": {"primary": "#1a73e8", "dark": "#202124", "accent": "#34a853"}
}
"""

import argparse
import json
import os
import sys
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


# --- Default Configuration ---
DEFAULT_COLORS = {
    "primary": "#1a73e8",
    "dark": "#202124",
    "accent": "#34a853",
    "light_gray": "#f8f9fa",
    "text_dark": "#202124",
    "text_light": "#5f6368",
    "white": "#ffffff"
}

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def get_text_color(hex_color):
    """Return white or dark text color based on background brightness."""
    r, g, b = hex_to_rgb(hex_color)
    brightness = (r * 299 + g * 587 + b * 114) / 1000
    return RGBColor(255, 255, 255) if brightness < 128 else RGBColor(32, 33, 36)


def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None):
    """Add a shape with optional fill and line colors."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(fill_color))
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = RGBColor(*hex_to_rgb(line_color))
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, 
                 font_color="#202124", bold=False, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*hex_to_rgb(font_color))
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_footer(slide, project_name, page_num, colors):
    """Add consistent footer with page number and project name."""
    # Bottom bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.2), 
              SLIDE_WIDTH, Inches(0.3), colors["light_gray"])
    
    # Page number
    add_text_box(slide, Inches(0.5), Inches(7.2), Inches(0.5), Inches(0.3),
                 str(page_num), font_size=10, font_color=colors["text_light"])
    
    # Project name
    add_text_box(slide, Inches(1), Inches(7.2), Inches(4), Inches(0.3),
                 project_name, font_size=10, font_color=colors["text_light"])


def create_cover(prs, data, colors):
    """Page 1 - Cover / Project Overview"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background gradient effect (using shapes)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
              SLIDE_WIDTH, SLIDE_HEIGHT, colors["primary"])
    
    # Project name - large centered
    add_text_box(slide, Inches(0.5), Inches(2.5), Inches(12.333), Inches(1),
                 data.get("project_name", "Project Name"),
                 font_size=48, font_color=colors["white"], bold=True, alignment=PP_ALIGN.CENTER)
    
    # Tagline
    add_text_box(slide, Inches(0.5), Inches(3.6), Inches(12.333), Inches(0.5),
                 data.get("tagline", ""),
                 font_size=24, font_color=colors["white"], alignment=PP_ALIGN.CENTER)
    
    # Value propositions
    value_props = data.get("value_props", [])
    y_pos = 4.5
    for i, prop in enumerate(value_props[:3]):
        add_text_box(slide, Inches(4), Inches(y_pos + i*0.5), Inches(5.333), Inches(0.4),
                     f"• {prop}", font_size=16, font_color=colors["white"])
    
    # Date
    date_str = datetime.now().strftime("%Y-%m")
    add_text_box(slide, Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5),
                 f"{date_str}", font_size=14, font_color=colors["white"], alignment=PP_ALIGN.CENTER)
    
    add_footer(slide, data.get("project_name", ""), 1, colors)
    return slide


def create_pain_points(prs, data, colors):
    """Page 2 - Market Pain Points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_text = "Market Pain Points" if data.get("language") == "en" else "市场痛点"
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8),
                 title_text, font_size=32, font_color=colors["dark"], bold=True)
    
    # Pain points cards
    pain_points = data.get("pain_points", [])
    card_width = Inches(3.8)
    card_height = Inches(2.5)
    start_x = 0.5
    gap = Inches(0.35)
    
    for i, pp in enumerate(pain_points[:3]):
        x = Inches(start_x + i * (3.8 + 0.35))
        # Card background
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), card_width, card_height,
                  colors["light_gray"])
        # Number
        add_text_box(slide, x, Inches(1.7), card_width, Inches(0.5),
                     str(i+1), font_size=36, font_color=colors["primary"], bold=True, alignment=PP_ALIGN.CENTER)
        # Title
        add_text_box(slide, x + Inches(0.2), Inches(2.3), card_width - Inches(0.4), Inches(0.5),
                     pp.get("title", f"Pain Point {i+1}"),
                     font_size=18, font_color=colors["dark"], bold=True)
        # Description
        add_text_box(slide, x + Inches(0.2), Inches(2.8), card_width - Inches(0.4), Inches(1),
                     pp.get("description", ""), font_size=14, font_color=colors["text_light"])
    
    # Market size section
    market = data.get("market_size", {})
    y_start = 4.5
    add_text_box(slide, Inches(0.5), Inches(y_start), Inches(4), Inches(0.5),
                 "Market Size" if data.get("language") == "en" else "市场规模",
                 font_size=24, font_color=colors["dark"], bold=True)
    
    stats = [
        ("TAM", market.get("tam", "[TBD]")),
        ("SAM", market.get("sam", "[TBD]")),
        ("SOM", market.get("som", "[TBD]"))
    ]
    
    for i, (label, value) in enumerate(stats):
        x = Inches(0.5 + i * 4)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(y_start + 0.7), Inches(3.5), Inches(1.2),
                  colors["primary"])
        add_text_box(slide, x, Inches(y_start + 0.9), Inches(3.5), Inches(0.5),
                     value, font_size=24, font_color=colors["white"], bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(y_start + 1.4), Inches(3.5), Inches(0.3),
                     label, font_size=14, font_color=colors["white"], alignment=PP_ALIGN.CENTER)
    
    add_footer(slide, data.get("project_name", ""), 2, colors)


def create_solution(prs, data, colors):
    """Page 3 - Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_text = "Solution" if data.get("language") == "en" else "解决方案"
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8),
                 title_text, font_size=32, font_color=colors["dark"], bold=True)
    
    # Features
    solution = data.get("solution", {})
    features = solution.get("features", [])
    
    for i, feature in enumerate(features[:4]):
        row = i // 2
        col = i % 2
        x = Inches(0.5 + col * 6.2)
        y = Inches(1.5 + row * 2.5)
        
        # Feature card
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.8), Inches(2.2),
                  colors["light_gray"])
        add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(5.2), Inches(0.5),
                     feature.get("title", f"Feature {i+1}"),
                     font_size=18, font_color=colors["primary"], bold=True)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.9), Inches(5.2), Inches(1),
                     feature.get("description", ""), font_size=14, font_color=colors["text_light"])
    
    add_footer(slide, data.get("project_name", ""), 3, colors)


def create_business_model(prs, data, colors):
    """Page 4 - Business Model"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_text = "Business Model" if data.get("language") == "en" else "商业模式"
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8),
                 title_text, font_size=32, font_color=colors["dark"], bold=True)
    
    model = data.get("business_model", {})
    
    # Revenue streams
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(0.5),
                 "Revenue Streams" if data.get("language") == "en" else "收入来源",
                 font_size=24, font_color=colors["dark"], bold=True)
    
    streams = model.get("revenue_streams", [])
    for i, stream in enumerate(streams[:3]):
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.2 + i*0.8), 
                  Inches(5.8), Inches(0.7), colors["primary"])
        add_text_box(slide, Inches(0.7), Inches(2.35 + i*0.8), Inches(5.4), Inches(0.4),
                     stream, font_size=14, font_color=colors["white"])
    
    # Pricing strategy
    add_text_box(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(0.5),
                 "Pricing Strategy" if data.get("language") == "en" else "定价策略",
                 font_size=24, font_color=colors["dark"], bold=True)
    
    pricing = model.get("pricing", "[Pricing strategy to be added]")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.2), 
              Inches(5.8), Inches(2.5), colors["light_gray"])
    add_text_box(slide, Inches(7), Inches(2.4), Inches(5.4), Inches(2),
                 pricing, font_size=14, font_color=colors["text_light"])
    
    add_footer(slide, data.get("project_name", ""), 4, colors)


def create_product_demo(prs, data, colors):
    """Page 5 - Product Demo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_text = "Product Demo" if data.get("language") == "en" else "产品展示"
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8),
                 title_text, font_size=32, font_color=colors["dark"], bold=True)
    
    # Central mockup placeholder
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(1.8), 
              Inches(6.333), Inches(4), colors["light_gray"])
    add_text_box(slide, Inches(3.5), Inches(3.5), Inches(6.333), Inches(0.5),
                 "[Insert Product Screenshot Here]" if data.get("language") == "en" else "[插入产品截图]",
                 font_size=16, font_color=colors["text_light"], alignment=PP_ALIGN.CENTER)
    
    # Product status
    status = data.get("product_status", {})
    status_text = []
    if status.get("demo"):
        status_text.append("✓ Demo Available")
    if status.get("users"):
        status_text.append(f"Users: {status['users']}")
    if status.get("metrics"):
        status_text.append(f"Metrics: {status['metrics']}")
    
    y_pos = 6
    for text in status_text:
        add_text_box(slide, Inches(0.5), Inches(y_pos), Inches(12.333), Inches(0.4),
                     text, font_size=14, font_color=colors["accent"])
        y_pos += 0.4
    
    add_footer(slide, data.get("project_name", ""), 5, colors)


def create_competitive_analysis(prs, data, colors):
    """Page 6 - Competitive Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_text = "Competitive Analysis" if data.get("language") == "en" else "竞品分析"
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8),
                 title_text, font_size=32, font_color=colors["dark"], bold=True)
    
    # Comparison table headers
    headers = ["Feature", "Us"] + [c.get("name", f"Comp {i+1}") for i, c in enumerate(data.get("competitors", [])[:3])]
    
    # Simplified table representation
    y_start = 1.5
    row_height = 0.6
    col_widths = [3, 2.5, 2.5, 2.5]
    
    # Header row
    x_pos = 0.5
    for i, header in enumerate(headers):
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(x_pos), Inches(y_start), 
                  Inches(col_widths[i]), Inches(row_height), colors["primary"])
        add_text_box(slide, Inches(x_pos), Inches(y_start + 0.15), 
                      Inches(col_widths[i]), Inches(0.3),
                      header, font_size=12, font_color=colors["white"], 
                      bold=True, alignment=PP_ALIGN.CENTER)
        x_pos += col_widths[i]
    
    # Sample feature rows
    features = ["Core Feature", "Pricing", "Ease of Use", "Support"]
    for j, feature in enumerate(features):
        y = y_start + (j + 1) * row_height
        x_pos = 0.5
        
        # Feature name
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(x_pos), Inches(y), 
                  Inches(col_widths[0]), Inches(row_height), colors["light_gray"])
        add_text_box(slide, Inches(x_pos + 0.1), Inches(y + 0.15), 
                      Inches(col_widths[0] - 0.2), Inches(0.3),
                      feature, font_size=11, font_color=colors["dark"])
        x_pos += col_widths[0]
        
        # Comparison cells
        for i in range(3):
            add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(x_pos), Inches(y), 
                      Inches(col_widths[i+1]), Inches(row_height), colors["white"])
            mark = "✓" if i == 0 else "○"
            add_text_box(slide, Inches(x_pos), Inches(y + 0.15), 
                          Inches(col_widths[i+1]), Inches(0.3),
                          mark, font_size=14, font_color=colors["accent"], alignment=PP_ALIGN.CENTER)
            x_pos += col_widths[i+1]
    
    add_footer(slide, data.get("project_name", ""), 6, colors)


def create_traction(prs, data, colors):
    """Page 7 - Traction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_text = "Traction" if data.get("language") == "en" else "项目进展"
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8),
                 title_text, font_size=32, font_color=colors["dark"], bold=True)
    
    # Metrics
    traction = data.get("traction", {})
    metrics = traction.get("metrics", "")
    
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), 
              Inches(12.333), Inches(1.2), colors["accent"])
    add_text_box(slide, Inches(0.7), Inches(1.8), Inches(11.933), Inches(0.6),
                 f"Key Metrics: {metrics}" if metrics else "Key metrics to be added",
                 font_size=18, font_color=colors["white"])
    
    # Timeline
    add_text_box(slide, Inches(0.5), Inches(3), Inches(12.333), Inches(0.5),
                 "Milestones" if data.get("language") == "en" else "里程碑",
                 font_size=24, font_color=colors["dark"], bold=True)
    
    milestones = traction.get("milestones", [])
    for i, milestone in enumerate(milestones[:4]):
        x = Inches(0.5 + i * 3.1)
        # Timeline dot
        add_shape(slide, MSO_SHAPE.OVAL, x + Inches(1.2), Inches(3.7), Inches(0.3), Inches(0.3), colors["primary"])
        # Line
        if i < 3:
            add_shape(slide, MSO_SHAPE.RECTANGLE, x + Inches(1.5), Inches(3.82), Inches(2.6), Inches(0.06), colors["primary"])
        # Text
        add_text_box(slide, x, Inches(4.2), Inches(2.8), Inches(0.8),
                     milestone, font_size=12, font_color=colors["text_light"], alignment=PP_ALIGN.CENTER)
    
    add_footer(slide, data.get("project_name", ""), 7, colors)


def create_roadmap(prs, data, colors):
    """Page 8 - Roadmap"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_text = "Roadmap" if data.get("language") == "en" else "产品路线图"
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8),
                 title_text, font_size=32, font_color=colors["dark"], bold=True)
    
    roadmap = data.get("roadmap", {})
    phases = [
        ("0-6 Months", roadmap.get("short", [])),
        ("6-18 Months", roadmap.get("mid", [])),
        ("18+ Months", roadmap.get("long", []))
    ]
    
    for i, (title, items) in enumerate(phases):
        x = Inches(0.5 + i * 4.2)
        # Phase card
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(3.9), Inches(5), colors["light_gray"])
        # Phase title
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(3.9), Inches(0.6), colors["primary"])
        add_text_box(slide, x, Inches(1.6), Inches(3.9), Inches(0.4),
                     title, font_size=16, font_color=colors["white"], bold=True, alignment=PP_ALIGN.CENTER)
        
        # Items
        for j, item in enumerate(items[:4]):
            add_text_box(slide, x + Inches(0.2), Inches(2.3 + j*0.6), Inches(3.5), Inches(0.5),
                         f"• {item}", font_size=12, font_color=colors["text_light"])
    
    add_footer(slide, data.get("project_name", ""), 8, colors)


def create_team(prs, data, colors):
    """Page 9 - Team"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_text = "Team" if data.get("language") == "en" else "团队介绍"
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8),
                 title_text, font_size=32, font_color=colors["dark"], bold=True)
    
    team = data.get("team", [])
    
    for i, member in enumerate(team[:4]):
        col = i % 4
        x = Inches(0.5 + col * 3.1)
        y = Inches(1.5)
        
        # Avatar placeholder
        add_shape(slide, MSO_SHAPE.OVAL, x + Inches(0.8), Inches(y), Inches(1.2), Inches(1.2), colors["light_gray"])
        
        # Name
        add_text_box(slide, x, Inches(y + 1.4), Inches(2.8), Inches(0.4),
                     member.get("name", f"Member {i+1}"),
                     font_size=16, font_color=colors["dark"], bold=True, alignment=PP_ALIGN.CENTER)
        
        # Role
        add_text_box(slide, x, Inches(y + 1.8), Inches(2.8), Inches(0.4),
                     member.get("role", ""),
                     font_size=12, font_color=colors["primary"], alignment=PP_ALIGN.CENTER)
        
        # Background
        add_text_box(slide, x + Inches(0.1), Inches(y + 2.3), Inches(2.6), Inches(1.5),
                     member.get("background", "")[:80] + "..." if len(member.get("background", "")) > 80 else member.get("background", ""),
                     font_size=11, font_color=colors["text_light"], alignment=PP_ALIGN.CENTER)
    
    add_footer(slide, data.get("project_name", ""), 9, colors)


def create_fundraising(prs, data, colors):
    """Page 10 - Fundraising"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_text = "Fundraising" if data.get("language") == "en" else "融资计划"
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8),
                 title_text, font_size=32, font_color=colors["dark"], bold=True)
    
    fundraising = data.get("fundraising", {})
    
    # Target amount
    amount = fundraising.get("amount", "[Amount TBD]")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), 
              Inches(6), Inches(1.5), colors["primary"])
    add_text_box(slide, Inches(0.7), Inches(1.9), Inches(5.6), Inches(0.6),
                 "Target Raise", font_size=14, font_color=colors["white"])
    add_text_box(slide, Inches(0.7), Inches(2.3), Inches(5.6), Inches(0.6),
                 amount, font_size=28, font_color=colors["white"], bold=True)
    
    # Equity
    equity = fundraising.get("equity", "[Equity TBD]")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.5), 
              Inches(5.833), Inches(1.5), colors["accent"])
    add_text_box(slide, Inches(7.2), Inches(1.9), Inches(5.433), Inches(0.6),
                 "Equity Offered", font_size=14, font_color=colors["white"])
    add_text_box(slide, Inches(7.2), Inches(2.3), Inches(5.433), Inches(0.6),
                 equity, font_size=28, font_color=colors["white"], bold=True)
    
    # Use of funds
    add_text_box(slide, Inches(0.5), Inches(3.5), Inches(12.333), Inches(0.5),
                 "Use of Funds" if data.get("language") == "en" else "资金用途",
                 font_size=24, font_color=colors["dark"], bold=True)
    
    use_of_funds = fundraising.get("use_of_funds", [])
    colors_list = [colors["primary"], colors["accent"], "#fbbc04", "#ea4335"]
    
    for i, fund in enumerate(use_of_funds[:4]):
        x = Inches(0.5 + i * 3.1)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.2), Inches(2.8), Inches(1.5),
                  colors_list[i % len(colors_list)])
        add_text_box(slide, x + Inches(0.2), Inches(4.8), Inches(2.4), Inches(0.6),
                     fund, font_size=14, font_color=colors["white"], alignment=PP_ALIGN.CENTER)
    
    # Contact info
    add_text_box(slide, Inches(0.5), Inches(6), Inches(12.333), Inches(0.5),
                 "Contact: [your@email.com]", font_size=14, font_color=colors["text_light"],
                 alignment=PP_ALIGN.CENTER)
    
    add_footer(slide, data.get("project_name", ""), 10, colors)


def main():
    parser = argparse.ArgumentParser(description="Generate a professional pitch deck")
    parser.add_argument("--input", "-i", required=True, help="Path to JSON input file")
    parser.add_argument("--output", "-o", help="Output PPTX file path (optional)")
    args = parser.parse_args()
    
    # Load input data
    if not os.path.exists(args.input):
        print(f"Error: Input file not found: {args.input}")
        sys.exit(1)
    
    with open(args.input, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    # Determine output filename
    if args.output:
        output_file = args.output
    else:
        project_name = data.get("project_name", "Startup").replace(" ", "_")
        output_file = f"{project_name}_BP.pptx"
    
    # Get colors from data or use defaults
    colors = DEFAULT_COLORS.copy()
    if "colors" in data:
        colors.update(data["colors"])
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # Generate all 10 slides
    create_cover(prs, data, colors)
    create_pain_points(prs, data, colors)
    create_solution(prs, data, colors)
    create_business_model(prs, data, colors)
    create_product_demo(prs, data, colors)
    create_competitive_analysis(prs, data, colors)
    create_traction(prs, data, colors)
    create_roadmap(prs, data, colors)
    create_team(prs, data, colors)
    create_fundraising(prs, data, colors)
    
    # Save
    prs.save(output_file)
    print(f"✅ Pitch deck generated: {os.path.abspath(output_file)}")
    print(f"   Slides: 10 | Format: 16:9 Widescreen")


if __name__ == "__main__":
    main()
